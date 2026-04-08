#!/usr/bin/env python3
"""Generate the 'Your Agent Control Center' workshop PPTX from the TBB Format template.

Usage: python generate_workshop.py

Clones slides from 'TBB Format.pptx' and replaces text in named shapes to produce
the workshop deck. All styling, colors, fonts, backgrounds, and decorative elements
come directly from the template.
"""
import copy
import re
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

TEMPLATE = Path(__file__).parent.parent / "TBB Format.pptx"
OUTPUT = Path(__file__).parent / "Your Agent Control Center - Workshop.pptx"

# Template slide indices (0-based)
TITLE_SLIDE = 0
SECTION_SLIDE = 1
CONTENT_SLIDE = 3
COMPARISON_SLIDE = 4
QUOTE_SLIDE = 5
PROCESS_SLIDE = 8
CLOSE_SLIDE = 12
TRANSITION_SLIDE = 19

# XML namespaces
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


def get_shape_name(sp_element):
    """Get the name of a shape from its XML element."""
    nvSpPr = sp_element.find(qn('p:nvSpPr'))
    if nvSpPr is None:
        return None
    cNvPr = nvSpPr.find(qn('p:cNvPr'))
    if cNvPr is None:
        return None
    return cNvPr.get('name')


def set_shape_text(sp_element, new_text):
    """Replace all text in a shape element while preserving formatting."""
    txBody = sp_element.find(qn('p:txBody'))
    if txBody is None:
        return

    paragraphs = txBody.findall(qn('a:p'))
    if not paragraphs:
        return

    # Set text in first paragraph's first run
    first_p = paragraphs[0]
    runs = first_p.findall(qn('a:r'))
    if runs:
        t = runs[0].find(qn('a:t'))
        if t is not None:
            t.text = new_text
        # Clear remaining runs
        for r in runs[1:]:
            t = r.find(qn('a:t'))
            if t is not None:
                t.text = ""
    else:
        # No runs, try direct text
        for t in first_p.findall(qn('a:t')):
            t.text = new_text
            break

    # Clear remaining paragraphs
    for p in paragraphs[1:]:
        for r in p.findall(qn('a:r')):
            t = r.find(qn('a:t'))
            if t is not None:
                t.text = ""


def set_shape_lines(sp_element, lines):
    """Replace text with multiple lines, preserving paragraph formatting."""
    txBody = sp_element.find(qn('p:txBody'))
    if txBody is None:
        return

    paragraphs = txBody.findall(qn('a:p'))
    if not paragraphs:
        return

    for i, line in enumerate(lines):
        if i < len(paragraphs):
            p = paragraphs[i]
            runs = p.findall(qn('a:r'))
            if runs:
                t = runs[0].find(qn('a:t'))
                if t is not None:
                    t.text = line
                for r in runs[1:]:
                    t = r.find(qn('a:t'))
                    if t is not None:
                        t.text = ""
        else:
            # Clone last paragraph and set text
            last_p = copy.deepcopy(paragraphs[-1])
            runs = last_p.findall(qn('a:r'))
            if runs:
                t = runs[0].find(qn('a:t'))
                if t is not None:
                    t.text = line
                for r in runs[1:]:
                    t = r.find(qn('a:t'))
                    if t is not None:
                        t.text = ""
            txBody.append(last_p)

    # Clear extra paragraphs
    for j in range(len(lines), len(paragraphs)):
        for r in paragraphs[j].findall(qn('a:r')):
            t = r.find(qn('a:t'))
            if t is not None:
                t.text = ""


def clone_slide(prs, template_idx):
    """Clone a template slide and return the new slide with editable XML."""
    template_slide = prs.slides[template_idx]
    slide_layout = template_slide.slide_layout

    new_slide = prs.slides.add_slide(slide_layout)

    # Remove all default shapes from new slide
    spTree = new_slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
    for child in list(spTree):
        tag = child.tag
        if tag == qn('p:nvGrpSpPr') or tag == qn('p:grpSpPr'):
            continue  # keep group shape properties
        spTree.remove(child)

    # Copy shapes from template
    template_spTree = template_slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
    for child in template_spTree:
        tag = child.tag
        if tag == qn('p:nvGrpSpPr') or tag == qn('p:grpSpPr'):
            continue
        spTree.append(copy.deepcopy(child))

    # Copy background
    template_cSld = template_slide._element.find(qn('p:cSld'))
    template_bg = template_cSld.find(qn('p:bg'))
    if template_bg is not None:
        new_cSld = new_slide._element.find(qn('p:cSld'))
        old_bg = new_cSld.find(qn('p:bg'))
        if old_bg is not None:
            new_cSld.remove(old_bg)
        new_cSld.insert(0, copy.deepcopy(template_bg))

    return new_slide


def edit_slide(slide, shape_texts):
    """Edit named shapes on a slide. shape_texts is {name: text_or_list}."""
    spTree = slide._element.find(qn('p:cSld')).find(qn('p:spTree'))

    for sp in spTree.iter(qn('p:sp')):
        name = get_shape_name(sp)
        if name and name in shape_texts:
            value = shape_texts[name]
            if isinstance(value, list):
                set_shape_lines(sp, value)
            else:
                set_shape_text(sp, value)


def delete_template_slides(prs, count):
    """Remove the first `count` slides from the presentation."""
    prs_element = prs.part._element
    sldIdLst = prs_element.find(qn('p:sldIdLst'))
    slide_ids = list(sldIdLst)
    for sid in slide_ids[:count]:
        sldIdLst.remove(sid)


# --- Workshop Content ---

SLIDES = [
    (TITLE_SLIDE, {
        "Overline": "THE BITCOIN BREAKDOWN PRESENTS",
        "MainTitle": "Your Agent Control Center",
        "Subtitle": "Stop Vibing. Start Directing.",
        "Meta": "WORKSHOP  \u2022  2026",
    }),
    (SECTION_SLIDE, {
        "SectionNum": "SECTION 01",
        "SectionTitle": "What We're Building",
        "SectionDesc": "An Agent Control Center on your computer -- not a chatbot, not a coding tool",
    }),
    (CONTENT_SLIDE, {
        "ContentOverline": "THE SETUP",
        "ContentTitle": "Three Things Working Together",
        "ContentBody": [
            "VS Code -- a free program that shows your files and folders",
            "Claude Code -- an AI agent that understands plain English",
            "Your folders -- real files on your real computer that you own",
            "You type what you want. Files appear. You own them forever.",
        ],
    }),
    (COMPARISON_SLIDE, {
        "TwoColTitle": "The Spectrum",
        "LeftHeader": "VIBING",
        "LeftContent": [
            "Chatting with AI in a browser",
            "Chat history trapped in the platform",
            "Gone when you close the tab",
            "Doesn't build over time",
        ],
        "RightHeader": "DIRECTING",
        "RightContent": [
            "Running an AI agent from your desktop",
            "Real files in your folders",
            "Everything remains when you close",
            "Workspace gets smarter over time",
        ],
    }),
    (CONTENT_SLIDE, {
        "ContentOverline": "THE SKILL",
        "ContentTitle": "You Already Know How to Do This",
        "ContentBody": [
            "The skill is talking clearly. Not coding. Not prompt engineering.",
            "If you can explain what you want to a coworker, you can direct an AI agent.",
            "You describe what you need. The agent builds it. You review. You redirect.",
            "The thinking is yours. The speed is the AI's.",
        ],
    }),
    (SECTION_SLIDE, {
        "SectionNum": "SECTION 02",
        "SectionTitle": "Let's Install It",
        "SectionDesc": "Everything you need -- about 10-15 minutes",
    }),
    (PROCESS_SLIDE, {
        "TimelineTitle": "Installation",
        "StepCircle0": "01", "StepTitle0": "VS CODE", "StepDesc0": "Go to code.visualstudio.com. Download. Install. Open it.",
        "StepCircle1": "02", "StepTitle1": "EXTENSION", "StepDesc1": "Extensions icon > Search 'Claude Code' > Install the one by Anthropic.",
        "StepCircle2": "03", "StepTitle2": "SIGN IN", "StepDesc2": "Click Claude icon > Sign in > Log in at claude.ai with Pro account ($20/mo).",
        "StepCircle3": "04", "StepTitle3": "FIRST FILE", "StepDesc3": "Create a folder. Open it in VS Code. Type a request. Watch the file appear.",
    }),
    (CONTENT_SLIDE, {
        "ContentOverline": "KEY CONCEPT",
        "ContentTitle": "Website vs. Agent",
        "ContentBody": [
            "claude.ai (the website) is a chatbot in a browser. That's vibing.",
            "Claude Code (in VS Code) is an agent on your computer. That's directing.",
            "Same company. Same AI. Completely different experience.",
            "You're paying for Pro because it unlocks Claude Code -- the agent.",
        ],
    }),
    (CONTENT_SLIDE, {
        "ContentOverline": "HANDS ON",
        "ContentTitle": "Your First Conversation",
        "ContentBody": [
            "Create a folder: 'My Agent Control Center' (Desktop or Documents)",
            "In VS Code: File > Open Folder > select it",
            "Type: 'Create a file called hello.md with an introduction to this workspace'",
            "Click Allow when asked. Watch the file appear. Click it. Read it.",
        ],
    }),
    (SECTION_SLIDE, {
        "SectionNum": "SECTION 03",
        "SectionTitle": "Now Do Something Real",
        "SectionDesc": "Pick a real project. Direct the agent. See the difference.",
    }),
    (CONTENT_SLIDE, {
        "ContentOverline": "YOUR PROJECT",
        "ContentTitle": "Pick Something You Actually Need",
        "ContentBody": [
            "A plan you haven't made. Research you've been putting off.",
            "A letter you haven't written. A budget for an upcoming event.",
            "A comparison you need to make. A guide you wish existed.",
            "Pick one. Something specific. Something you'll actually use.",
        ],
    }),
    (QUOTE_SLIDE, {
        "QuoteText": "\"I'm planning a monthly meetup. I need a logistics plan, three session agendas, and a budget assuming I'm paying out of pocket. Put everything in organized files.\"",
        "QuoteAttrib": "\u2014 EXAMPLE PROMPT",
    }),
    (CONTENT_SLIDE, {
        "ContentOverline": "THE WORKFLOW",
        "ContentTitle": "Direct the Revision",
        "ContentBody": [
            "Read what Claude created. Some will be right. Some won't.",
            "Don't just accept it. Find what's off. Be specific about changes.",
            "The file updates on your computer -- not a new chat message.",
            "Direct. Review. Redirect. Files get better with each pass.",
        ],
    }),
    (SECTION_SLIDE, {
        "SectionNum": "SECTION 04",
        "SectionTitle": "Give It Memory",
        "SectionDesc": "CLAUDE.md -- a plain text file that makes every conversation smarter",
    }),
    (CONTENT_SLIDE, {
        "ContentOverline": "CONTEXT ENGINEERING",
        "ContentTitle": "CLAUDE.md: Your Agent's Memory",
        "ContentBody": [
            "Create a file called CLAUDE.md in your project folder.",
            "The agent reads it automatically at the start of every conversation.",
            "Tell it: who you are, what this project is, how you like to work.",
            "Next time you open the folder, the agent already knows all of this.",
        ],
    }),
    (CONTENT_SLIDE, {
        "ContentOverline": "THE LEVERAGE POINT",
        "ContentTitle": "Specificity Changes Everything",
        "ContentBody": [
            "'I'm a teacher' = generic results you can't use.",
            "'I teach AP Bio to smart but bored 17-year-olds, and I need",
            "2-3 page study guides with quizzes' = results you can use.",
            "The more specific your CLAUDE.md, the better every conversation gets.",
        ],
    }),
    (TRANSITION_SLIDE, {
        "TransTitle": "Three Rules",
        "TransSubtitle": "The difference between decent results and results you can use",
    }),
    (PROCESS_SLIDE, {
        "TimelineTitle": "Three Rules",
        "StepCircle0": "01", "StepTitle0": "VERIFY", "StepDesc0": "Don't just trust polished output. Read it. Check facts. Run the numbers.",
        "StepCircle1": "02", "StepTitle1": "ONE AT A TIME", "StepDesc1": "One project, one task. Check the output. Then the next thing.",
        "StepCircle2": "03", "StepTitle2": "THINK FIRST", "StepDesc2": "What do I need? What does 'done' look like? What constraints matter?",
        "StepCircle3": "", "StepTitle3": "", "StepDesc3": "",
    }),
    (CONTENT_SLIDE, {
        "ContentOverline": "TAKEAWAY",
        "ContentTitle": "What You Have Now",
        "ContentBody": [
            "VS Code + Claude Code installed and working on your computer",
            "A project folder with real files -- files you own, not chat logs",
            "A CLAUDE.md that makes every future conversation smarter",
            "A workspace that accumulates knowledge over time",
        ],
    }),
    (CLOSE_SLIDE, {
        "ClosingHead": "Stop Vibing.\nStart Directing.",
        "ClosingSub": "Take home the best practices handout. Keep experimenting.",
        "ContactInfo": "thebitcoinbreakdown.com",
    }),
]


def main():
    prs = Presentation(str(TEMPLATE))
    template_count = len(prs.slides)

    for template_idx, shape_texts in SLIDES:
        slide = clone_slide(prs, template_idx)
        edit_slide(slide, shape_texts)

    delete_template_slides(prs, template_count)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Workshop slides: {len(SLIDES)}")


if __name__ == "__main__":
    main()
